from __future__ import annotations

import argparse
import ctypes
import hashlib
import html
import json
import os
import subprocess
import sys
import threading
import time
import urllib.parse
import urllib.request
import webbrowser
from datetime import datetime
from http.server import SimpleHTTPRequestHandler, ThreadingHTTPServer
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont

W, H = 640, 420


def default_desktop() -> Path:
    return Path(os.environ.get("USERPROFILE", str(Path.home()))) / "Desktop"


def default_data_dir() -> Path:
    base = Path(os.environ.get("LOCALAPPDATA", str(Path.home())))
    return base / "DesktopVisualGallery"


def font(size: int, bold: bool = False):
    fonts = Path(os.environ.get("WINDIR", r"C:\Windows")) / "Fonts"
    for name in (["msyhbd.ttc", "arialbd.ttf"] if bold else ["msyh.ttc", "arial.ttf"]):
        path = fonts / name
        if path.exists():
            return ImageFont.truetype(str(path), size)
    return ImageFont.load_default()


def fitted(image: Image.Image) -> Image.Image:
    image = image.convert("RGB")
    image.thumbnail((W - 28, H - 28), Image.Resampling.LANCZOS)
    output = Image.new("RGB", (W, H), "#f4f6f9")
    output.paste(image, ((W - image.width) // 2, (H - image.height) // 2))
    return output


def type_card(ext: str, name: str, lines: str = "") -> Image.Image:
    colors = {".xlsx": "#1f8f5f", ".xls": "#1f8f5f", ".csv": "#1f8f5f", ".pdf": "#dc4c3e", ".docx": "#3678c9", ".doc": "#3678c9", ".pptx": "#d96735", ".zip": "#7b61a8", ".mp4": "#3a4759"}
    image = Image.new("RGB", (W, H), "#f4f6f9")
    draw = ImageDraw.Draw(image)
    color = colors.get(ext, "#596579")
    draw.rounded_rectangle((38, 34, W - 38, H - 34), radius=24, fill="white", outline="#d9dee7", width=2)
    draw.rounded_rectangle((70, 68, 210, 208), radius=20, fill=color)
    draw.text((140, 138), (ext[1:] or "FILE").upper()[:5], font=font(30, True), fill="white", anchor="mm")
    short = name[:29] + "…" if len(name) > 30 else name
    draw.text((240, 88), short, font=font(22), fill="#182230")
    for index, line in enumerate(lines.splitlines()[:7]):
        draw.text((240, 135 + index * 31), line[:34], font=font(18), fill="#667085")
    return image


def render_preview(path: Path) -> Image.Image:
    ext = path.suffix.lower()
    try:
        if ext in {".png", ".jpg", ".jpeg", ".bmp", ".gif", ".webp", ".tif", ".tiff"}:
            with Image.open(path) as image:
                return fitted(image.copy())
        if ext == ".pdf":
            import pypdfium2 as pdfium
            pdf = pdfium.PdfDocument(str(path))
            return fitted(pdf[0].render(scale=1.3).to_pil()) if len(pdf) else type_card(ext, path.stem)
        if ext in {".xlsx", ".xlsm"}:
            from openpyxl import load_workbook
            book = load_workbook(path, read_only=True, data_only=True)
            sheet = next((s for s in book.worksheets if s.max_row and s.max_column), book.worksheets[0])
            rows = [["" if value is None else str(value)[:16] for value in row] for row in sheet.iter_rows(min_row=1, max_row=min(sheet.max_row, 14), max_col=min(sheet.max_column, 7), values_only=True)]
            image = Image.new("RGB", (W, H), "white")
            draw = ImageDraw.Draw(image)
            draw.rectangle((0, 0, W, 48), fill="#1f8f5f")
            draw.text((20, 24), sheet.title[:26], font=font(22), fill="white", anchor="lm")
            columns = max([len(row) for row in rows] or [1])
            column_width = (W - 24) / columns
            for row_index, row in enumerate(rows):
                for column_index in range(columns):
                    x, y = 12 + column_index * column_width, 58 + row_index * 25
                    draw.rectangle((x, y, x + column_width, y + 25), fill="#e9f5ef" if row_index == 0 else "white", outline="#cfd6df")
                    value = row[column_index] if column_index < len(row) else ""
                    draw.text((x + 4, y + 4), value[:14], font=font(13), fill="#273142")
            book.close()
            return image
        if ext == ".docx":
            from docx import Document
            lines = [p.text.strip() for p in Document(path).paragraphs if p.text.strip()]
            return type_card(ext, path.stem, "\n".join(lines[:7]))
        if ext == ".pptx":
            from pptx import Presentation
            deck = Presentation(path)
            lines = []
            if deck.slides:
                lines = [shape.text.strip() for shape in deck.slides[0].shapes if getattr(shape, "has_text_frame", False) and shape.text.strip()]
            return type_card(ext, path.stem, "\n".join(lines[:7]))
    except Exception:
        pass
    return type_card(ext, path.stem, "Open the original file")


def file_type(ext: str) -> str:
    if ext in {".xlsx", ".xls", ".xlsm", ".csv"}: return "Spreadsheet"
    if ext == ".pdf": return "PDF"
    if ext in {".docx", ".doc", ".rtf"}: return "Document"
    if ext in {".pptx", ".ppt"}: return "Presentation"
    if ext in {".png", ".jpg", ".jpeg", ".bmp", ".gif", ".webp", ".tif", ".tiff"}: return "Image"
    if ext in {".mp4", ".mov", ".avi", ".mkv"}: return "Video"
    if ext in {".zip", ".rar", ".7z"}: return "Archive"
    return "Other"


def build_index(root: Path, data_dir: Path) -> int:
    thumbs = data_dir / "thumbs"
    thumbs.mkdir(parents=True, exist_ok=True)
    files = sorted((p for p in root.rglob("*") if p.is_file() and p.suffix.lower() != ".lnk" and p.name != "desktop.ini" and not p.name.startswith("~$")), key=lambda p: p.stat().st_mtime, reverse=True)
    records = []
    for path in files:
        stat = path.stat()
        key = hashlib.sha1(str(path).encode("utf-8")).hexdigest()[:16] + ".jpg"
        thumbnail = thumbs / key
        if not thumbnail.exists() or thumbnail.stat().st_mtime < stat.st_mtime:
            render_preview(path).save(thumbnail, "JPEG", quality=82, optimize=True)
        relative = path.relative_to(root)
        folder = str(relative.parent) if relative.parent != Path(".") else root.name
        ext = path.suffix.lower()
        records.append({"name": path.name, "path": str(path), "thumb": "thumbs/" + key, "folder": folder, "type": file_type(ext), "size": stat.st_size, "sizeText": f"{stat.st_size / 1048576:.1f} MB" if stat.st_size >= 1048576 else f"{stat.st_size / 1024:.0f} KB", "modified": datetime.fromtimestamp(stat.st_mtime).strftime("%Y-%m-%d %H:%M"), "mtime": stat.st_mtime})
    (data_dir / "index.html").write_text(PAGE.replace("__DATA__", json.dumps(records, ensure_ascii=False).replace("</", "<\\/")), encoding="utf-8")
    (data_dir / "manifest.json").write_text(json.dumps(records, ensure_ascii=False, indent=2), encoding="utf-8")
    return len(records)


def copy_text(text: str) -> None:
    GMEM_MOVEABLE, CF_UNICODETEXT = 0x0002, 13
    data = (text + "\0").encode("utf-16-le")
    kernel32, user32 = ctypes.windll.kernel32, ctypes.windll.user32
    kernel32.GlobalAlloc.argtypes, kernel32.GlobalAlloc.restype = [ctypes.c_uint, ctypes.c_size_t], ctypes.c_void_p
    kernel32.GlobalLock.argtypes, kernel32.GlobalLock.restype = [ctypes.c_void_p], ctypes.c_void_p
    kernel32.GlobalUnlock.argtypes, kernel32.GlobalFree.argtypes = [ctypes.c_void_p], [ctypes.c_void_p]
    user32.OpenClipboard.argtypes, user32.OpenClipboard.restype = [ctypes.c_void_p], ctypes.c_int
    user32.SetClipboardData.argtypes, user32.SetClipboardData.restype = [ctypes.c_uint, ctypes.c_void_p], ctypes.c_void_p
    handle = kernel32.GlobalAlloc(GMEM_MOVEABLE, len(data))
    pointer = kernel32.GlobalLock(handle)
    if not handle or not pointer: raise OSError("Clipboard memory allocation failed")
    ctypes.memmove(pointer, data, len(data)); kernel32.GlobalUnlock(handle)
    for _ in range(12):
        if user32.OpenClipboard(None): break
        time.sleep(0.05)
    else:
        kernel32.GlobalFree(handle); raise OSError("Clipboard is busy")
    try:
        if not user32.EmptyClipboard() or not user32.SetClipboardData(CF_UNICODETEXT, handle): raise OSError("Clipboard write failed")
        handle = None
    finally:
        user32.CloseClipboard()
        if handle: kernel32.GlobalFree(handle)


def reveal_file(path: Path) -> None:
    pidl, attrs = ctypes.c_void_p(), ctypes.c_ulong()
    shell32, ole32 = ctypes.windll.shell32, ctypes.windll.ole32
    init = ole32.CoInitialize(None)
    try:
        result = shell32.SHParseDisplayName(str(path), None, ctypes.byref(pidl), 0, ctypes.byref(attrs))
        if result != 0 or not pidl.value: raise OSError(f"SHParseDisplayName failed: 0x{result & 0xFFFFFFFF:08X}")
        try:
            result = shell32.SHOpenFolderAndSelectItems(pidl, 0, None, 0)
            if result != 0: raise OSError(f"SHOpenFolderAndSelectItems failed: 0x{result & 0xFFFFFFFF:08X}")
        finally: ole32.CoTaskMemFree(pidl)
    finally:
        if init in (0, 1): ole32.CoUninitialize()


def handler_factory(root: Path, data_dir: Path):
    class Handler(SimpleHTTPRequestHandler):
        def __init__(self, *args, **kwargs): super().__init__(*args, directory=str(data_dir), **kwargs)
        def log_message(self, *_): pass
        def do_GET(self):
            parsed = urllib.parse.urlparse(self.path)
            if parsed.path in {"/open", "/reveal", "/copy"}:
                raw = urllib.parse.parse_qs(parsed.query).get("path", [""])[0]
                try:
                    target = Path(raw).resolve(strict=True); target.relative_to(root)
                    if not target.is_file(): raise ValueError
                except (OSError, ValueError):
                    self.send_error(403, "Only files under the indexed root are allowed"); return
                try:
                    if parsed.path == "/open": os.startfile(str(target))
                    elif parsed.path == "/reveal": reveal_file(target)
                    else: copy_text(str(target))
                    self.send_response(204); self.end_headers()
                except OSError as exc: self.send_error(500, str(exc))
                return
            if parsed.path == "/": self.path = "/index.html"
            super().do_GET()
    return Handler


def main():
    parser = argparse.ArgumentParser(description="Local-first visual file gallery for Windows")
    parser.add_argument("--root", type=Path, default=default_desktop())
    parser.add_argument("--data-dir", type=Path, default=default_data_dir())
    parser.add_argument("--port", type=int, default=8765)
    parser.add_argument("--no-browser", action="store_true")
    parser.add_argument("--refresh-only", action="store_true")
    args = parser.parse_args()
    root, data_dir = args.root.resolve(), args.data_dir.resolve()
    if not root.is_dir(): raise SystemExit(f"Indexed root does not exist: {root}")
    count = build_index(root, data_dir)
    print(f"Indexed {count} files into {data_dir}")
    if args.refresh_only: return
    url = f"http://127.0.0.1:{args.port}/"
    try:
        with urllib.request.urlopen(url, timeout=0.35) as response:
            if response.status == 200:
                if not args.no_browser: webbrowser.open(url)
                return
    except Exception: pass
    server = ThreadingHTTPServer(("127.0.0.1", args.port), handler_factory(root, data_dir))
    if not args.no_browser: threading.Timer(0.7, lambda: webbrowser.open(url)).start()
    server.serve_forever()


PAGE = r'''<!doctype html><html lang="en"><head><meta charset="utf-8"><meta name="viewport" content="width=device-width,initial-scale=1"><title>Desktop Visual Gallery</title><style>
:root{--bg:#f5f7fb;--ink:#162033;--muted:#667085;--line:#e2e7ef;--blue:#356ae6}*{box-sizing:border-box}body{margin:0;font-family:"Segoe UI",sans-serif;background:var(--bg);color:var(--ink)}header{position:sticky;top:0;z-index:5;background:rgba(245,247,251,.95);backdrop-filter:blur(12px);border-bottom:1px solid var(--line)}.head{max-width:1500px;margin:auto;padding:22px 28px 14px}.top{display:flex;align-items:end;gap:18px;justify-content:space-between}h1{font-size:27px;margin:0 0 4px}.sub{color:var(--muted);font-size:14px}.search{width:min(520px,45vw);padding:12px 16px;border:1px solid #ccd4e1;border-radius:12px;background:white;font-size:15px}.bar{display:flex;gap:8px;flex-wrap:wrap;margin-top:15px}.chip{border:1px solid #d4dbe6;background:white;border-radius:999px;padding:8px 13px;cursor:pointer;color:#475467}.chip.active{color:white;background:var(--blue);border-color:var(--blue)}main{max-width:1500px;margin:auto;padding:24px 28px 50px}.meta{display:flex;justify-content:space-between;color:var(--muted);font-size:14px;margin-bottom:15px}.grid{display:grid;grid-template-columns:repeat(auto-fill,minmax(270px,1fr));gap:17px}.card{background:white;border:1px solid var(--line);border-radius:15px;overflow:hidden;box-shadow:0 3px 12px rgba(16,24,40,.04);transition:.16s}.card:hover{transform:translateY(-3px);box-shadow:0 10px 24px rgba(16,24,40,.12);border-color:#b8c6e0}.preview{display:block;height:180px;background:#eef1f5;overflow:hidden;border-bottom:1px solid var(--line)}.preview img{width:100%;height:100%;object-fit:contain;display:block}.body{padding:13px 14px 14px;min-height:150px}.name{font-weight:650;line-height:1.4;height:2.8em;overflow:hidden;word-break:break-word}.info{display:flex;justify-content:space-between;color:var(--muted);font-size:12px;margin-top:10px}.tag{color:var(--blue);font-weight:650}.note{font-size:12px;color:#7a8596;margin-top:7px;white-space:nowrap;overflow:hidden;text-overflow:ellipsis}.actions{display:flex;gap:7px;margin-top:11px}.actions a{font-size:12px;text-decoration:none;color:#315fce;background:#eef3ff;border-radius:7px;padding:6px 8px}.sort{border:1px solid #d4dbe6;background:white;border-radius:9px;padding:7px 10px;color:#475467}@media(max-width:720px){.top{display:block}.search{width:100%;margin-top:14px}.head,main{padding-left:16px;padding-right:16px}.meta{display:block}.meta>span:last-child{display:block;margin-top:10px}.meta select{width:100%;margin-top:7px}.grid{grid-template-columns:1fr;gap:12px}}
</style></head><body><header><div class="head"><div class="top"><div><h1>Desktop Visual Gallery</h1><div class="sub">See it first, remember it later. Open the original file from any card.</div></div><input id="q" class="search" placeholder="Search by filename, folder, type, or date"></div><div class="bar" id="filters"></div></div></header><main><div class="meta"><span id="count"></span><span><select id="folder" class="sort"><option value="All folders">All folders</option></select> <select id="sort" class="sort"><option value="new">Recently modified</option><option value="old">Oldest modified</option><option value="name">Name</option><option value="large">Largest</option></select></span></div><div id="grid" class="grid"></div></main><iframe name="action-target" hidden></iframe><script>
const DATA=__DATA__,TYPES=['All','Spreadsheet','PDF','Document','Presentation','Image','Video','Archive','Other'];let active='All';const esc=s=>s.replace(/[&<>"']/g,c=>({'&':'&amp;','<':'&lt;','>':'&gt;','"':'&quot;',"'":'&#39;'}[c]));function render(){let q=document.querySelector('#q').value.trim().toLowerCase(),f=document.querySelector('#folder').value;let a=DATA.filter(x=>(active==='All'||x.type===active)&&(f==='All folders'||x.folder===f)&&(!q||(x.name+' '+x.type+' '+x.modified+' '+x.folder).toLowerCase().includes(q)));let s=document.querySelector('#sort').value;a.sort(s==='old'?(x,y)=>x.mtime-y.mtime:s==='name'?(x,y)=>x.name.localeCompare(y.name):s==='large'?(x,y)=>y.size-x.size:(x,y)=>y.mtime-x.mtime);document.querySelector('#count').textContent=`Showing ${a.length} of ${DATA.length} files`;document.querySelector('#grid').innerHTML=a.map(x=>`<div class="card" title="${esc(x.path)}"><a class="preview" href="/open?path=${encodeURIComponent(x.path)}" target="action-target"><img loading="lazy" src="${x.thumb}" alt=""></a><div class="body"><div class="name">${esc(x.name)}</div><div class="info"><span class="tag">${x.type}</span><span>${x.sizeText}</span></div><div class="note">${esc(x.folder)} · ${x.modified}</div><div class="actions"><a href="/open?path=${encodeURIComponent(x.path)}" target="action-target">Open</a><a href="/reveal?path=${encodeURIComponent(x.path)}" target="action-target">Show file</a><a href="/copy?path=${encodeURIComponent(x.path)}" target="action-target" onclick="const t=this;t.textContent='Copied';setTimeout(()=>t.textContent='Copy path',1200)">Copy path</a></div></div></div>`).join('')}const folders=[...new Set(DATA.map(x=>x.folder))].sort();document.querySelector('#folder').innerHTML+=folders.map(f=>`<option value="${esc(f)}">${esc(f)}</option>`).join('');document.querySelector('#filters').innerHTML=TYPES.map(t=>`<button class="chip ${t==='All'?'active':''}" data-t="${t}">${t}</button>`).join('');document.querySelector('#filters').onclick=e=>{if(!e.target.dataset.t)return;active=e.target.dataset.t;document.querySelectorAll('.chip').forEach(x=>x.classList.toggle('active',x.dataset.t===active));render()};document.querySelector('#q').oninput=render;document.querySelector('#sort').onchange=render;document.querySelector('#folder').onchange=render;render();
</script></body></html>'''


if __name__ == "__main__":
    main()
